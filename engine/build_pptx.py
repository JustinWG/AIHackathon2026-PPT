"""
A3 — PPTX assembler (Person A).

Public API (do not change signatures without team sync):
    build_pptx(spec, master_path, out_path) -> str

Supports two layouts (both must exist in the master):
    - "Team19 Title Slide"        : title (idx 0) + subtitle (idx 1)
    - "Team19 Title and Content"  : title (idx 0) + content (idx 1) -> bullets OR table

Speaker notes always go into the slide's notes, never on the slide itself.
"""
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor


NOTES_TEMPLATE = """\
Aim: {aim}

Time: {time}

Instructions: {instructions}

Reflective question: {reflective_q}

Debrief: {debrief}"""


class AssemblyError(Exception):
    pass


# Layouts whose filled text should be forced to white, regardless of the
# theme/layout default. Useful when a branded layout has a dark background
# but its placeholder text would otherwise inherit a dark color.
WHITE_FONT_LAYOUTS = {"Team19 Title Slide"}


def _load_layout_map(layouts_path: str = "schema/layouts.json") -> dict:
    data = json.loads(Path(layouts_path).read_text())
    block_map = data.get("block_to_layout", {})
    if not block_map or any("FILL_ME" in v for v in block_map.values()):
        raise AssemblyError(
            "schema/layouts.json block_to_layout is not filled in yet. "
            "Run engine/layouts_report.py first."
        )
    return block_map


def _get_layout(prs: Presentation, layout_name: str):
    # Search every slide master, not just the first. Decks imported from
    # Google Slides often have many masters; prs.slide_layouts only exposes
    # the first master's layouts, so a needed layout can be "missing" while
    # actually living under another master.
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == layout_name:
                return layout
    available = sorted(
        {l.name for m in prs.slide_masters for l in m.slide_layouts}
    )
    raise AssemblyError(
        f"Layout '{layout_name}' not found in master. Available: {available}"
    )


def _placeholder_by_idx(slide, idx: int):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _force_white(text_frame) -> None:
    """Set every run in the text frame to white (overrides theme/layout color)."""
    for para in text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = _WHITE


def _fill_title(slide, title: str, white: bool = False) -> None:
    ph = _placeholder_by_idx(slide, 0)
    if ph is not None:
        ph.text = title
        if white:
            _force_white(ph.text_frame)


def _fill_bullets(content_ph, bullets: list[str]) -> None:
    tf = content_ph.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            tf.paragraphs[0].text = bullet
        else:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0


def _add_table(slide, table_data: dict, content_ph=None) -> None:
    headers = table_data["headers"]
    rows = table_data["rows"]
    cols = len(headers)
    total_rows = len(rows) + 1

    if content_ph is not None:
        left, top, width = content_ph.left, content_ph.top, content_ph.width
        height = min(content_ph.height, int(Inches(0.45 * total_rows)))
    else:
        left, top = Inches(0.5), Inches(2.0)
        width, height = Inches(9.0), Inches(0.45 * total_rows)

    tbl = slide.shapes.add_table(total_rows, cols, left, top, width, height).table

    for col_idx, header in enumerate(headers):
        tbl.cell(0, col_idx).text = str(header)
    for row_idx, row in enumerate(rows):
        for col_idx, cell_val in enumerate(row):
            tbl.cell(row_idx + 1, col_idx).text = str(cell_val)


def _fill_content(slide, bullets: list[str], table_data, white: bool = False) -> None:
    """Fill the content placeholder (idx 1) with a table if present, else bullets."""
    content_ph = _placeholder_by_idx(slide, 1)
    if table_data:
        # Remove the empty content placeholder so the table isn't layered over it.
        if content_ph is not None:
            ref = content_ph
            _add_table(slide, table_data, content_ph=ref)
            ref._element.getparent().remove(ref._element)
        else:
            _add_table(slide, table_data)
    elif bullets and content_ph is not None:
        _fill_bullets(content_ph, bullets)
        if white:
            _force_white(content_ph.text_frame)


def _write_notes(slide, notes: dict) -> None:
    slide.notes_slide.notes_text_frame.text = NOTES_TEMPLATE.format(**notes)


def _remove_all_slides(prs: Presentation) -> None:
    """Fully drop template slides: both the sldId entries and the slide parts.

    Removing only the sldIdLst entries leaves orphan slide parts in the package,
    which makes the writer emit duplicate part names (and can trigger a repair
    prompt in PowerPoint). Dropping the parts too keeps the package clean.
    """
    sld_id_lst = prs.slides._sldIdLst
    part = prs.part
    for sld_id in list(sld_id_lst):
        rId = sld_id.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        sld_id_lst.remove(sld_id)
        if rId and rId in part.rels:
            part.drop_rel(rId)


def build_pptx(spec: dict, master_path: str, out_path: str) -> str:
    """
    Fills the master's layouts from spec, writes speaker notes, returns out_path.
    Raises AssemblyError on bad layout name or malformed spec.
    """
    block_map = _load_layout_map()
    prs = Presentation(master_path)

    # Drop any slides shipped with the template; keep the layouts.
    _remove_all_slides(prs)

    for slide_spec in spec["slides"]:
        block = slide_spec["block"]
        layout_name = slide_spec.get("layout") or block_map.get(block)
        if not layout_name or layout_name == "FILL_AFTER_LAYOUT_REPORT":
            layout_name = block_map.get(block)
        if not layout_name:
            raise AssemblyError(f"No layout for block '{block}'")

        layout = _get_layout(prs, layout_name)
        slide = prs.slides.add_slide(layout)

        white_font = layout_name in WHITE_FONT_LAYOUTS
        _fill_title(slide, slide_spec["title"], white=white_font)
        _fill_content(
            slide,
            slide_spec.get("bullets") or [],
            slide_spec.get("table"),
            white=white_font,
        )
        _write_notes(slide, slide_spec["notes"])

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return out_path


if __name__ == "__main__":
    import sys
    spec_path   = sys.argv[1] if len(sys.argv) > 1 else "schema/sample_spec.json"
    master_path = sys.argv[2] if len(sys.argv) > 2 else "master/Maverx - Presentation Style Guide for Hackaton.pptx"
    out_path    = sys.argv[3] if len(sys.argv) > 3 else "output/test.pptx"

    spec = json.loads(Path(spec_path).read_text())
    result = build_pptx(spec, master_path, out_path)
    print(f"Written: {result}")
