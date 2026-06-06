"""
A1 — Layout report (Person A, run this first).

Usage:
    python engine/layouts_report.py master/maverx_master.pptx

Prints every layout name + index + placeholder details, and writes
schema/layouts.json with a `layouts` array for B to reference.
Then manually fill in `block_to_layout` in layouts.json.
"""
import json
import sys
from pathlib import Path
from pptx import Presentation
from pptx.enum.text import PP_ALIGN


def report(master_path: str, out_path: str = "schema/layouts.json") -> None:
    prs = Presentation(master_path)
    layouts = []

    for idx, layout in enumerate(prs.slide_layouts):
        placeholders = []
        for ph in layout.placeholders:
            placeholders.append({
                "idx":  ph.placeholder_format.idx,
                "type": str(ph.placeholder_format.type),
                "name": ph.name,
            })
        entry = {"index": idx, "name": layout.name, "placeholders": placeholders}
        layouts.append(entry)
        print(f"[{idx:2d}] {layout.name}")
        for ph in placeholders:
            print(f"       idx={ph['idx']}  type={ph['type']}  name={ph['name']}")

    existing = json.loads(Path(out_path).read_text()) if Path(out_path).exists() else {}
    existing["layouts"] = layouts
    Path(out_path).write_text(json.dumps(existing, indent=2))
    print(f"\nWritten to {out_path}")
    print("→ Now fill in block_to_layout manually using the layout names above.")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python engine/layouts_report.py master/maverx_master.pptx")
        sys.exit(1)
    report(sys.argv[1])
